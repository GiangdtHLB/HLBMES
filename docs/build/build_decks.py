#!/usr/bin/env python3
"""Dựng 3 bộ slide PowerPoint training cho MES Bia Hạ Long.
Thiết kế: nền sáng cho nội dung, nền navy cho bìa/divider; motif = vòng tròn
màu chứa số/chữ; bảng màu navy + teal + amber (chủ đề nhà máy bia).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Bảng màu ----
NAVY = RGBColor(0x1E, 0x27, 0x61)
NAVY2 = RGBColor(0x16, 0x1D, 0x49)
TEAL = RGBColor(0x02, 0x80, 0x90)
SEA = RGBColor(0x00, 0xA8, 0x96)
AMBER = RGBColor(0xF2, 0xA9, 0x00)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x22, 0x2A, 0x3A)
GREY = RGBColor(0x5A, 0x63, 0x72)
CARD = RGBColor(0xEF, 0xF3, 0xF9)
CARD2 = RGBColor(0xE7, 0xF0, 0xF2)
LINEC = RGBColor(0xD5, 0xDD, 0xE8)

SW, SH = 13.333, 7.5
BODY = "Calibri"
HEAD = "Cambria"


def new_prs():
    p = Presentation()
    p.slide_width = Inches(SW)
    p.slide_height = Inches(SH)
    return p


def blank(prs, bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = bg
    sh.line.fill.background()
    sh.shadow.inherit = False
    return s


def rect(s, l, t, w, h, fill, rounded=False, line=None, lw=1.0):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def circle(s, l, t, d, fill, label="", lc=WHITE, ls=16):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.fill.background(); c.shadow.inherit = False
    if label:
        tf = c.text_frame; tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.size = Pt(ls); r.font.bold = True; r.font.color.rgb = lc; r.font.name = BODY
    return c


def tb(s, l, t, w, h, paras, anchor=MSO_ANCHOR.TOP):
    """paras: list of dict{runs:[{t,s,c,b,i,font}], align, sa, sb, line}"""
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    tf.vertical_anchor = anchor
    for i, pa in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pa.get("align", PP_ALIGN.LEFT)
        if "sa" in pa: p.space_after = Pt(pa["sa"])
        if "sb" in pa: p.space_before = Pt(pa["sb"])
        p.line_spacing = pa.get("line", 1.0)
        for rn in pa["runs"]:
            r = p.add_run(); r.text = rn["t"]
            r.font.size = Pt(rn.get("s", 15))
            r.font.bold = rn.get("b", False)
            r.font.italic = rn.get("i", False)
            r.font.name = rn.get("font", BODY)
            r.font.color.rgb = rn.get("c", INK)
    return box


def one(t, s=15, c=INK, b=False, i=False, font=BODY, align=PP_ALIGN.LEFT, sa=0, line=1.0):
    return {"runs": [{"t": t, "s": s, "c": c, "b": b, "i": i, "font": font}],
            "align": align, "sa": sa, "line": line}


def heading(s, title, kicker=""):
    if kicker:
        tb(s, 0.6, 0.42, 12, 0.32, [one(kicker.upper(), 12, TEAL, b=True)])
    tb(s, 0.6, 0.72, 12.1, 0.9, [one(title, 30, NAVY, b=True, font=HEAD)])


def title_slide(prs, kicker, title, subtitle):
    s = blank(prs, NAVY)
    rect(s, 0, 0, SW, SH, NAVY)
    circle(s, 11.0, -1.4, 3.4, NAVY2)
    circle(s, -1.2, 5.4, 3.2, NAVY2)
    tb(s, 0.9, 1.5, 11, 0.5, [one("🍺 MES BIA HẠ LONG · NHÀ MÁY ĐÔNG MAI", 15, AMBER, b=True)])
    tb(s, 0.9, 2.15, 11.6, 2.0, [one(title, 46, WHITE, b=True, font=HEAD, line=1.02)])
    tb(s, 0.92, 4.35, 11, 0.6, [one(subtitle, 19, ICE)])
    rect(s, 0.95, 5.25, 2.6, 0.06, TEAL)
    tb(s, 0.9, 6.6, 11, 0.4,
       [one(kicker + "  ·  Phiên bản 0.1.0-mvp  ·  Tài liệu hướng dẫn nội bộ", 12, ICE)])
    return s


def divider(prs, num, title, sub=""):
    s = blank(prs, NAVY)
    circle(s, 0.9, 2.55, 1.7, TEAL, num, WHITE, 40)
    tb(s, 2.95, 2.5, 9.5, 1.2, [one(title, 34, WHITE, b=True, font=HEAD)])
    if sub:
        tb(s, 2.98, 3.7, 9.3, 0.8, [one(sub, 16, ICE)])
    return s


def card(s, l, t, w, h, badge, header, lines, accent=TEAL, hsize=15, bsize=12.5):
    rect(s, l, t, w, h, CARD, rounded=True, line=LINEC, lw=0.75)
    d = 0.5
    circle(s, l + 0.22, t + 0.22, d, accent, badge, WHITE, 14)
    tb(s, l + 0.85, t + 0.18, w - 1.0, 0.55,
       [one(header, hsize, NAVY, b=True)], anchor=MSO_ANCHOR.MIDDLE)
    paras = [one(x, bsize, INK if not x.startswith("•") else GREY, line=1.06, sa=2) for x in lines]
    tb(s, l + 0.28, t + 0.82, w - 0.5, h - 0.95, paras)


def footer(s, txt):
    tb(s, 0.6, 7.04, 12.1, 0.32,
       [one(txt, 9.5, GREY, align=PP_ALIGN.LEFT)])


def simple_table(s, l, t, w, rows, col_w, header_fill=NAVY, hsize=11, bsize=10.5, rh=0.34):
    """rows[0] = header. col_w fractions sum≈1."""
    x = l
    widths = [w * f for f in col_w]
    th = 0.46
    # header
    cx = l
    for ci, htext in enumerate(rows[0]):
        rect(s, cx, t, widths[ci], th, header_fill, line=WHITE, lw=0.5)
        tb(s, cx + 0.08, t, widths[ci] - 0.14, th,
           [one(htext, hsize, WHITE, b=True)], anchor=MSO_ANCHOR.MIDDLE)
        cx += widths[ci]
    # body
    y = t + th
    for ri, row in enumerate(rows[1:]):
        fill = WHITE if ri % 2 == 0 else CARD
        cx = l
        for ci, val in enumerate(row):
            rect(s, cx, y, widths[ci], rh, fill, line=LINEC, lw=0.5)
            tb(s, cx + 0.08, y, widths[ci] - 0.14, rh,
               [one(val, bsize, INK)], anchor=MSO_ANCHOR.MIDDLE)
            cx += widths[ci]
        y += rh
    return y


def stat(s, l, t, w, big, label, accent=TEAL):
    rect(s, l, t, w, 1.7, CARD, rounded=True, line=LINEC, lw=0.75)
    tb(s, l, t + 0.22, w, 0.9, [one(big, 40, accent, b=True, align=PP_ALIGN.CENTER, font=HEAD)])
    tb(s, l + 0.1, t + 1.12, w - 0.2, 0.5, [one(label, 12, GREY, align=PP_ALIGN.CENTER, line=1.0)])


# =====================================================================
# DECK 1 — KIẾN TRÚC
# =====================================================================
def deck_architecture():
    prs = new_prs()
    title_slide(prs, "Tài liệu Kiến trúc Chuẩn",
                "Tài liệu Kiến trúc Chuẩn",
                "Hệ thống Điều hành Sản xuất (MES) cho nhà máy bia")

    # 2. Tổng quan
    s = blank(prs); heading(s, "MES làm gì?", "Tổng quan")
    tb(s, 0.6, 1.75, 12.1, 0.8,
       [one("Điều phối toàn bộ vòng đời sản xuất bia — từ lệnh đến thành phẩm — "
            "trên một nền tảng thống nhất, có truy xuất và kiểm toán đầy đủ.", 16, INK, line=1.15)])
    flow = ["Lệnh SX", "Điều độ (WO)", "Mẻ (Batch)", "QC release", "Phả hệ", "Audit"]
    x = 0.6
    for i, f in enumerate(flow):
        rect(s, x, 2.8, 1.78, 0.7, NAVY if i % 2 == 0 else TEAL, rounded=True)
        tb(s, x, 2.8, 1.78, 0.7, [one(f, 12.5, WHITE, b=True, align=PP_ALIGN.CENTER)],
           anchor=MSO_ANCHOR.MIDDLE)
        if i < len(flow) - 1:
            tb(s, x + 1.78, 2.8, 0.22, 0.7, [one("›", 20, AMBER, b=True, align=PP_ALIGN.CENTER)],
               anchor=MSO_ANCHOR.MIDDLE)
        x += 2.0
    card(s, 0.6, 3.95, 3.86, 2.6, "1", "Phân hệ lõi",
         ["Lệnh & điều độ", "Công thức & BOM", "Mẻ & ISA-88", "Chất lượng & QC Lab",
          "Truy xuất & EBR"], NAVY)
    card(s, 4.73, 3.95, 3.86, 2.6, "2", "Phân hệ vận hành",
         ["Kho NVL & WMS", "OEE & dừng máy", "Bảo trì & kiểm định",
          "Năng lượng", "Bao bì tuần hoàn"], TEAL)
    card(s, 8.86, 3.95, 3.86, 2.6, "3", "Nền tảng số",
         ["Historian / Edge OT", "Lập lịch tối ưu", "Trợ lý AI (tư vấn)",
          "Cổng API tích hợp", "Audit hash-chain"], SEA)
    footer(s, "13/13 phân hệ MES hardcore đã hoàn thành")

    # 3. Nguyên tắc
    s = blank(prs); heading(s, "Nguyên tắc kiến trúc", "Triết lý thiết kế")
    items = [("M", "Modular monolith", "Một tiến trình, chia theo bounded context — dễ tách microservice sau"),
             ("B", "Bounded context", "Mỗi module sở hữu dữ liệu riêng; truy cập chéo qua service layer"),
             ("R", "REST / JSON", "OpenAPI tự sinh tại /docs; mọi API theo chuẩn REST"),
             ("D", "RDBMS chuẩn", "SQLite (dev) ↔ PostgreSQL 16 (prod), đổi qua biến môi trường"),
             ("I", "SoR bất biến", "Snapshot công thức · audit append-only · EBR khóa được"),
             ("H", "Human-in-the-loop", "AI chỉ tư vấn; mọi hành động cần con người + đúng quyền")]
    for i, (bd, hd, ds) in enumerate(items):
        col = i % 2; row = i // 2
        l = 0.6 + col * 6.15; t = 1.75 + row * 1.62
        rect(s, l, t, 5.95, 1.45, CARD, rounded=True, line=LINEC, lw=0.75)
        circle(s, l + 0.25, t + 0.42, 0.62, [NAVY, TEAL, SEA, AMBER, TEAL, NAVY][i], bd, WHITE, 20)
        tb(s, l + 1.1, t + 0.18, 4.7, 0.5, [one(hd, 15.5, NAVY, b=True)])
        tb(s, l + 1.1, t + 0.66, 4.7, 0.7, [one(ds, 11.5, GREY, line=1.08)])

    # 4. Sơ đồ tầng
    s = blank(prs); heading(s, "Kiến trúc phân tầng", "Layered architecture")
    layers = [("CLIENT", "Web UI · Kiosk · Swagger · phần mềm ngoài (API) · edge OPC UA", ICE, NAVY),
              ("MIDDLEWARE", "request-id · rate-limit · log · /metrics · ánh xạ lỗi → HTTP", SEA, WHITE),
              ("ROUTERS", "~31 router · ~180 REST endpoint", TEAL, WHITE),
              ("SERVICES", "23 module quy tắc nghiệp vụ (batches, recipes, quality, bom…)", TEAL, WHITE),
              ("MODELS (ORM)", "~40 bảng theo bounded context", NAVY, WHITE),
              ("DATABASE", "SQLite (dev)  ·  PostgreSQL 16 + Alembic (prod)", NAVY2, WHITE)]
    y = 1.7
    for name, desc, fill, fg in layers:
        rect(s, 0.6, y, 8.7, 0.78, fill, rounded=True)
        tb(s, 0.85, y, 2.5, 0.78, [one(name, 13.5, fg, b=True)], anchor=MSO_ANCHOR.MIDDLE)
        tb(s, 3.2, y, 6.0, 0.78, [one(desc, 11, fg, line=1.0)], anchor=MSO_ANCHOR.MIDDLE)
        y += 0.88
    rect(s, 9.6, 1.7, 3.1, 5.06, CARD, rounded=True, line=LINEC, lw=0.75)
    tb(s, 9.8, 1.9, 2.7, 0.5, [one("Xuyên suốt", 13, AMBER, b=True)])
    tb(s, 9.8, 2.45, 2.75, 4.2,
       [one("security.py", 12, NAVY, b=True, sa=1), one("RBAC · SoD · scope", 10.5, GREY, sa=8),
        one("audit.py", 12, NAVY, b=True, sa=1), one("hash-chain bất biến", 10.5, GREY, sa=8),
        one("config · logging", 12, NAVY, b=True, sa=1), one("Settings · request-id", 10.5, GREY, sa=8),
        one("ratelimit · metrics", 12, NAVY, b=True, sa=1), one("Redis · Prometheus", 10.5, GREY)])
    footer(s, "Quy tắc phụ thuộc: Router → Service → Model (service không gọi ngược router)")

    # 5. Tech stack
    s = blank(prs); heading(s, "Ngăn xếp công nghệ", "Tech stack")
    rows = [["Lớp", "Công nghệ", "Ghi chú"],
            ["Web framework", "FastAPI 0.111", "OpenAPI tự sinh, lifespan, middleware"],
            ["ORM", "SQLAlchemy 2.0", "DeclarativeBase, session-per-request"],
            ["Validation/cấu hình", "Pydantic 2 + settings", "Gom mọi biến MES_*"],
            ["CSDL", "SQLite / PostgreSQL 16", "Đổi qua MES_DATABASE_URL"],
            ["Migration", "Alembic", "6+ migration; prod: upgrade head"],
            ["Edge OT", "asyncua (OPC UA)", "Tag chuẩn Weihenstephan → historian"],
            ["AI (tùy chọn)", "anthropic / Claude", "claude-opus-4-8; offline → engine luật"],
            ["Frontend", "Vanilla JS + SVG", "Zero-build, chạy offline"],
            ["Đóng gói / CI", "Docker · GitHub Actions", "app+Postgres; ruff+pytest+build"]]
    simple_table(s, 0.6, 1.75, 12.1, rows, [0.26, 0.30, 0.44], rh=0.42, bsize=11)

    # 6. Mô hình dữ liệu
    s = blank(prs); heading(s, "Mô hình dữ liệu (~40 bảng)", "12 bounded context")
    ctxs = [("Auth", "User, Session"), ("Master", "Product, Material, Line"),
            ("Lệnh & điều độ", "PO, WorkOrder, ScheduleSlot"), ("Công thức", "Recipe, Version, Change"),
            ("Thực thi mẻ", "Batch, PhaseRun, Reading, OEE"), ("Vật tư & phả hệ", "Lot, GenealogyEdge, Dispense"),
            ("Chất lượng", "QC Result, Deviation, CAPA, Sample"), ("Nấu-Lọc-Chiết", "Brew, Ferment, Filter, Bottle"),
            ("Kho & bao bì", "StockMovement, Pallet, Case"), ("CMMS & năng lượng", "Equipment, Incident, Energy"),
            ("Historian", "HistorianPoint (UNS tag)"), ("Audit & AI", "AuditLog, Signature, Conversation")]
    for i, (h, d) in enumerate(ctxs):
        col = i % 3; row = i // 3
        l = 0.6 + col * 4.08; t = 1.75 + row * 1.28
        rect(s, l, t, 3.9, 1.12, CARD, rounded=True, line=LINEC, lw=0.75)
        tb(s, l + 0.2, t + 0.13, 3.6, 0.4, [one(h, 13, NAVY, b=True)])
        tb(s, l + 0.2, t + 0.55, 3.6, 0.5, [one(d, 10, GREY, line=1.04)])

    # 7. State machines
    s = blank(prs); heading(s, "Máy trạng thái cốt lõi", "Sai transition → HTTP 409")
    sms = [("Mẻ (Batch)", "planned → ready → running → held → completed → closed", NAVY),
           ("Công thức", "draft → review → approved → effective → suspended/obsolete", TEAL),
           ("Lệnh SX (WO)", "planned → released → in_progress → completed → closed", SEA),
           ("Deviation", "open → triage → investigation → disposition → approval → closed", AMBER),
           ("Chất lượng", "pending → on_hold ⇄ released → rejected   (FAIL → on_hold)", TEAL),
           ("Phase ISA-88", "idle → running → held → complete  (+ aborted)", NAVY)]
    y = 1.8
    for name, flow, col in sms:
        rect(s, 0.6, y, 3.0, 0.74, col, rounded=True)
        tb(s, 0.6, y, 3.0, 0.74, [one(name, 13, WHITE, b=True, align=PP_ALIGN.CENTER)],
           anchor=MSO_ANCHOR.MIDDLE)
        rect(s, 3.75, y, 8.95, 0.74, CARD, rounded=True, line=LINEC, lw=0.75)
        tb(s, 3.95, y, 8.6, 0.74, [one(flow, 12.5, INK, font="Consolas")], anchor=MSO_ANCHOR.MIDDLE)
        y += 0.86
    footer(s, "Chỉ phiên bản công thức 'effective' mới được dùng để chạy mẻ")

    # 8. Luồng nghiệp vụ một mẻ
    s = blank(prs); heading(s, "Vòng đời một mẻ sản xuất", "End-to-end")
    steps = ["PO → WorkOrder → Dispatch sinh Batch (snapshot recipe bất biến)",
             "Kiểm tra tồn BOM ↔ tồn khả dụng (thiếu → chặn 409)",
             "ready → running · Consume NVL theo FEFO + genealogy edge",
             "Ghi actual tham số · ghi QC (PASS/FAIL số học; FAIL → on_hold)",
             "Produce lô output · QA release (chặn nếu còn FAIL)",
             "completed → closed · EBR: ký điện tử → khóa (content_hash)"]
    for i, st in enumerate(steps):
        t = 1.75 + i * 0.86
        circle(s, 0.6, t, 0.6, [NAVY, TEAL, SEA, TEAL, SEA, AMBER][i], str(i + 1), WHITE, 18)
        rect(s, 1.4, t, 11.3, 0.66, CARD if i % 2 == 0 else CARD2, rounded=True, line=LINEC, lw=0.6)
        tb(s, 1.65, t, 10.9, 0.66, [one(st, 13, INK, line=1.0)], anchor=MSO_ANCHOR.MIDDLE)

    # 9. Bảo mật
    s = blank(prs); heading(s, "Bảo mật & phân quyền", "Defense in depth")
    card(s, 0.6, 1.75, 5.95, 2.3, "1", "Xác thực",
         ["Mật khẩu PBKDF2 100k vòng", "Token Bearer · phiên 12h",
          "Buộc đổi mật khẩu mặc định", "Fallback header chỉ ở dev"], NAVY)
    card(s, 6.75, 1.75, 5.95, 2.3, "2", "RBAC + quyền thao tác",
         ["5 vai trò: operator/supervisor/", "   qa/engineer/admin",
          "19 quyền (require_perm)", "Admin: toàn quyền"], TEAL)
    card(s, 0.6, 4.2, 5.95, 2.3, "3", "SoD + data-scoping",
         ["Soạn ≠ duyệt; ghi QC ≠ release", "Ký ≠ khóa EBR",
          "Scope theo line / khu vực / loại test"], SEA)
    card(s, 6.75, 4.2, 5.95, 2.3, "4", "Toàn vẹn (21 CFR Part 11)",
         ["Audit append-only + hash-chain", "verify-chain phát hiện giả mạo",
          "E-signature re-auth + lý do", "EBR khóa → mẻ bất biến"], AMBER)

    # 10. Tích hợp
    s = blank(prs); heading(s, "Tích hợp & Edge OT", "Mở rộng ra ngoài")
    card(s, 0.6, 1.75, 3.86, 4.6, "A", "Cổng API mở /api/v1",
         ["X-API-Key scope read/write", "batches · inventory · OEE",
          "energy · quality · trace", "Feed sự kiện (since_seq)",
          "Nhận event ngoài (ghi qua", "   record_audit, không gãy", "   hash-chain)"], TEAL)
    card(s, 4.73, 1.75, 3.86, 4.6, "B", "Edge & Historian",
         ["Tag UNS brewery/site/area/", "   device/metric", "ingest qua X-API-Key write",
          "downsample min/avg/max", "edge_sim mô phỏng gateway",
          "opcua_edge: client OPC UA", "   thật (Weihenstephan)"], SEA)
    card(s, 8.86, 1.75, 3.86, 4.6, "C", "Barcode / Kiosk",
         ["Tem Code39 + QR (segno)", "/api/scan phân giải mã",
          "lô / mẻ / WO / đơn hàng", "Kiosk cảm ứng cho xưởng",
          "cấp liệu nhanh + in tem"], NAVY)

    # 11. AI
    s = blank(prs); heading(s, "Lớp AI — chỉ tư vấn", "Human-in-the-loop")
    tb(s, 0.6, 1.7, 12.1, 0.6,
       [one("Mọi tool AI đều read-only — không tool nào đổi setpoint, điều khiển thiết bị "
            "hay ghi dữ liệu sản xuất.", 15, INK, line=1.1, b=False)])
    card(s, 0.6, 2.5, 3.86, 3.7, "✦", "AI vận hành",
         ["GET /api/ai/insights", "Quét dữ liệu thật →",
          "cảnh báo & đề xuất có", "mức ưu tiên (tồn/hạn,",
          "kiểm định, sự cố, QC,", "OEE thấp)"], TEAL, hsize=15)
    card(s, 4.73, 2.5, 3.86, 3.7, "✦", "Trợ lý chat",
         ["Claude claude-opus-4-8", "(tool-use) hoặc engine",
          "luật offline", "Bộ nhớ hội thoại server",
          "Streaming SSE", "8 tool MES read-only"], SEA, hsize=15)
    card(s, 8.86, 2.5, 3.86, 3.7, "✦", "Kiểm soát chi phí",
         ["Rate-limit 20/phút/phiên", "Hạn mức 300 chat/ngày",
          "Log model/token/USD/", "   latency mỗi lượt",
          "Tác vụ nền: ai_report,", "   recall"], NAVY, hsize=15)

    # 12. Vận hành & triển khai
    s = blank(prs); heading(s, "Vận hành & triển khai", "Production-grade")
    rows = [["Hạng mục", "Hiện thực"],
            ["Health / readiness", "GET /api/health — kiểm tra kết nối DB + version"],
            ["Metrics", "GET /metrics (Prometheus): http, AI, rate-limit, audit-chain"],
            ["Logging", "Structured + request-id (X-Request-ID) + log chi phí AI"],
            ["Rate-limit", "in-process sliding-window hoặc Redis (tự fallback)"],
            ["Đóng gói", "Docker compose: app FastAPI + PostgreSQL 16"],
            ["Migration", "Alembic upgrade head (prod) · create_all (dev)"],
            ["Backup", "backup.sh / restore.sh + test_restore (verify hash-chain)"],
            ["CI / Test", "GitHub Actions: ruff + pytest (30/30) + docker build"]]
    simple_table(s, 0.6, 1.75, 12.1, rows, [0.26, 0.74], rh=0.46, bsize=11)

    # 13. Giới hạn
    s = blank(prs); heading(s, "Giới hạn & ranh giới tích hợp thật", "Chủ ý cho MVP")
    lims = ["Tích hợp thiết bị OT chạy ở dạng mô phỏng edge + client OPC UA demo — production cần trỏ PLC/SCADA thật",
            "Historian dùng SQLite; production swap TimescaleDB/InfluxDB (điểm cắm đã sẵn)",
            "Chưa có SSO/OIDC + MFA (cần IdP); HTTPS qua reverse proxy (mẫu nginx sẵn)",
            "availability là kiểm tra tư vấn, chưa giữ chỗ tồn (TOCTOU nếu tạo mẻ đồng thời)",
            "Worker nền in-process; quy mô rất lớn → Celery/RQ + Redis broker",
            "Hồ sơ CSV/IQ-OQ-PQ & UAT theo ca thật thuộc quy trình tại site"]
    for i, lx in enumerate(lims):
        t = 1.8 + i * 0.82
        circle(s, 0.6, t, 0.5, AMBER, "!", NAVY, 16)
        rect(s, 1.3, t, 11.4, 0.66, CARD, rounded=True, line=LINEC, lw=0.6)
        tb(s, 1.55, t, 11.0, 0.66, [one(lx, 12, INK, line=1.0)], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, "Mọi phần mô phỏng đều có điểm tích hợp chuẩn để cắm thiết bị/DB thật về sau")

    prs.save("docs/MES-KienTruc.pptx")
    print("✓ docs/MES-KienTruc.pptx", len(prs.slides.__iter__.__self__._sldIdLst))


# =====================================================================
# DECK 2 — TÍNH NĂNG
# =====================================================================
def deck_features():
    prs = new_prs()
    title_slide(prs, "Danh sách Tính năng",
                "Danh sách Tính năng",
                "Toàn bộ năng lực phần mềm theo phân hệ")

    # 2. Stats
    s = blank(prs); heading(s, "Phạm vi phần mềm", "Một cái nhìn")
    stat(s, 0.6, 1.9, 2.85, "~31", "router REST", NAVY)
    stat(s, 3.65, 1.9, 2.85, "~180", "endpoint API", TEAL)
    stat(s, 6.7, 1.9, 2.85, "~40", "bảng dữ liệu", SEA)
    stat(s, 9.75, 1.9, 2.85, "28", "tab giao diện", AMBER)
    stat(s, 0.6, 3.85, 2.85, "23", "service nghiệp vụ", TEAL)
    stat(s, 3.65, 3.85, 2.85, "13/13", "phân hệ hardcore", NAVY)
    stat(s, 6.7, 3.85, 2.85, "30/30", "test tự động", SEA)
    stat(s, 9.75, 3.85, 2.85, "5", "vai trò + 19 quyền", AMBER)
    tb(s, 0.6, 5.95, 12.1, 0.8,
       [one("Mỗi tính năng đều có endpoint API và quyền yêu cầu rõ ràng — chi tiết trong tài liệu Danh sách Tính năng.",
            14, GREY, i=True, line=1.1)])

    # 3. Sản xuất
    s = blank(prs); heading(s, "Phân hệ sản xuất", "Từ lệnh đến mẻ")
    card(s, 0.6, 1.75, 3.86, 4.7, "1", "Lệnh & điều độ",
         ["Tạo lệnh ERP (PO)", "Work Order ngày/ca/line", "Dispatch phát mẻ",
          "planned vs actual", "Lập lịch Gantt tối ưu", "CIP + né bảo trì"], NAVY)
    card(s, 4.73, 1.75, 3.86, 4.7, "2", "Công thức & BOM",
         ["Version + workflow duyệt", "BOM scale theo mẻ", "Dung sai ± %",
          "Kiểm tra tồn trước mẻ", "NVL thay thế", "Yield theo công đoạn",
          "Change-control e-sign"], TEAL)
    card(s, 8.86, 1.75, 3.86, 4.7, "3", "Mẻ & ISA-88",
         ["Snapshot recipe bất biến", "State machine", "Consume / produce",
          "Đường cong lên men", "Đối chiếu BOM↔thực tế", "Phase nấu/lên men/lọc/CIP"], SEA)

    # 4. Chất lượng
    s = blank(prs); heading(s, "Chất lượng & phòng Lab", "QA / QC / LIMS")
    card(s, 0.6, 1.75, 5.95, 4.6, "Q", "Chất lượng cơ bản",
         ["Ghi kết quả QC (PASS/FAIL số học)", "FAIL → tự động ON HOLD",
          "Hold / Release (quyền QA)", "Release chặn nếu còn FAIL chưa đóng",
          "Deviation workflow open → closed"], TEAL, bsize=13)
    card(s, 6.75, 1.75, 5.95, 4.6, "L", "QC Lab nâng cao",
         ["SPC control chart I-MR (UCL/LCL)", "Luật Western Electric · Cp/Cpk",
          "CAPA: khắc phục / phòng ngừa", "COA — phiếu phân tích cho mẻ",
          "LIMS-lite: đăng ký → test → xong"], NAVY, bsize=13)

    # 5. Kho & logistics
    s = blank(prs); heading(s, "Kho & Logistics", "NVL · thành phẩm · bao bì")
    card(s, 0.6, 1.75, 3.86, 4.6, "K", "Kho NVL",
         ["Nhập / xuất / hoàn / sang ngang", "Sổ cái bất biến",
          "Tồn on-hand", "Thẻ kho (số dư lũy kế)", "Hạn sử dụng (FEFO)",
          "BC nhập-xuất-tồn"], NAVY, bsize=12.5)
    card(s, 4.73, 1.75, 3.86, 4.6, "W", "Kho TP (WMS)",
         ["Đóng pallet (+ case + barcode)", "Vị trí kho (bin/cold/dock)",
          "Putaway / ship", "% lấp đầy", "Phân giải barcode",
          "Tem Code39"], TEAL, bsize=12.5)
    card(s, 8.86, 1.75, 3.86, 4.6, "B", "Bao bì tuần hoàn",
         ["Vỏ chai / két-gông / keg", "Tồn + lưu hành ngoài",
          "Tiền cược (deposit)", "Biến động: nhập/xuất/", "   thu hồi/loại bỏ/kiểm kê",
          "Chặn kiểm kê âm"], SEA, bsize=12.5)

    # 6. Thiết bị & hiệu suất
    s = blank(prs); heading(s, "Thiết bị & hiệu suất", "OEE · CMMS · năng lượng")
    card(s, 0.6, 1.75, 3.86, 4.6, "O", "OEE & dừng máy",
         ["OEE = A × P × Q", "Gauge + phân rã A/P/Q", "Cây lý do dừng máy",
          "Pareto + % tích lũy", "6 big losses", "MTBF / MTTR"], NAVY, bsize=12.5)
    card(s, 4.73, 1.75, 3.86, 4.6, "M", "Bảo trì & kiểm định",
         ["Thiết bị + phụ tùng (tồn min)", "Sự cố + xử lý",
          "Kế hoạch (tự đánh dấu quá hạn)", "Kiểm định valid/due/overdue",
          "Phóng xạ, van an toàn,", "   hiệu chuẩn TBĐ"], TEAL, bsize=12.5)
    card(s, 8.86, 1.75, 3.86, 4.6, "E", "Năng lượng",
         ["Nhóm: điện/nước/hơi/khí", "Số đọc ngày (upsert)",
          "Biểu đồ ngày theo nhóm", "Tổng hợp tháng",
          "Theo khu vực"], SEA, bsize=12.5)

    # 7. Truy xuất, EBR, audit
    s = blank(prs); heading(s, "Truy xuất, hồ sơ & toàn vẹn", "Compliance")
    card(s, 0.6, 1.75, 3.86, 4.6, "T", "Truy xuất & Recall",
         ["Truy ngược TP → NVL", "Truy xuôi NVL → TP",
          "Recall simulation", "Đồ thị genealogy có hướng",
          "Đếm lô ảnh hưởng + ms"], TEAL, bsize=12.5)
    card(s, 4.73, 1.75, 3.86, 4.6, "E", "EBR & e-signature",
         ["Dossier step-by-step", "Timeline thao tác từ audit",
          "BOM / QC / deviation / men", "Ký điện tử re-auth (Part 11)",
          "Khóa → mẻ bất biến"], NAVY, bsize=12.5)
    card(s, 8.86, 1.75, 3.86, 4.6, "A", "Audit bất biến",
         ["Append-only (không sửa/xóa)", "hash-chain tamper-evident",
          "verify-chain phát hiện", "   giả mạo", "seq UNIQUE chống race"], SEA, bsize=12.5)

    # 8. AI & tích hợp
    s = blank(prs); heading(s, "AI, tích hợp & nền tảng số", "")
    card(s, 0.6, 1.75, 3.86, 4.6, "✦", "Trợ lý AI",
         ["Chat có bộ nhớ + streaming", "AI insights (cảnh báo ưu tiên)",
          "8 tool MES read-only", "Manifest cho AI agent/MCP",
          "Tác vụ nền (job)"], TEAL, bsize=12.5, hsize=15)
    card(s, 4.73, 1.75, 3.86, 4.6, "↔", "Cổng tích hợp",
         ["API mở /api/v1 (X-API-Key)", "Feed + nhận sự kiện",
          "Quản lý API key (scope)", "Webhook (HMAC secret)"], NAVY, bsize=12.5, hsize=15)
    card(s, 8.86, 1.75, 3.86, 4.6, "📡", "Edge & Realtime",
         ["Historian time-series", "Tag UNS + downsample",
          "Tab Realtime auto 4s", "edge_sim + OPC UA thật"], SEA, bsize=12.5, hsize=14)

    # 9. Highlight 13/13
    s = blank(prs, NAVY)
    tb(s, 0.8, 0.8, 11.7, 0.9, [one("13 / 13 phân hệ MES hardcore — hoàn thành", 30, WHITE, b=True, font=HEAD)])
    subs = ["Lệnh & điều độ", "EBR điện tử", "Recipe/BOM nâng cao", "Tích hợp thiết bị",
            "Historian/time-series", "Material consumption", "Quality hardcore (SPC)",
            "OEE/downtime", "Barcode/RFID/kiosk", "Phân quyền sâu",
            "E-sign & audit bất biến", "Kiến trúc production", "ISA-88 / LIMS / WMS"]
    for i, x in enumerate(subs):
        col = i % 2; row = i // 2
        l = 0.8 + col * 6.0; t = 1.95 + row * 0.74
        circle(s, l, t, 0.46, SEA if i % 2 == 0 else AMBER, "✓", NAVY, 13)
        tb(s, l + 0.6, t, 5.2, 0.46, [one(x, 14, WHITE)], anchor=MSO_ANCHOR.MIDDLE)

    prs.save("docs/MES-TinhNang.pptx")
    print("✓ docs/MES-TinhNang.pptx")


# =====================================================================
# DECK 3 — HƯỚNG DẪN SỬ DỤNG
# =====================================================================
def deck_userguide():
    prs = new_prs()
    title_slide(prs, "Sách Hướng dẫn Sử dụng",
                "Hướng dẫn Sử dụng",
                "Dành cho tất cả tài khoản · tài liệu đào tạo nội bộ")

    # 2. Mục tiêu
    s = blank(prs); heading(s, "Tài liệu này dành cho ai?", "Bắt đầu")
    tb(s, 0.6, 1.7, 12.1, 0.7,
       [one("Mọi người dùng MES — đọc phần chung trước, rồi tìm phần dành cho chức danh của bạn.",
            16, INK, line=1.15)])
    card(s, 0.6, 2.6, 3.86, 3.6, "A", "Phần chung",
         ["Đăng nhập", "Đổi mật khẩu", "Giao diện & menu", "Hiểu phân quyền"], NAVY)
    card(s, 4.73, 2.6, 3.86, 3.6, "C", "Theo chức danh",
         ["10 tài khoản", "Menu thấy được", "Quyền & phạm vi", "Thao tác từng bước"], TEAL)
    card(s, 8.86, 2.6, 3.86, 3.6, "D", "Quy trình & FAQ",
         ["End-to-end 1 mẻ", "Kiosk xưởng", "Trợ lý AI", "Xử lý sự cố"], SEA)

    # 3. Đăng nhập
    s = blank(prs); heading(s, "Đăng nhập & phiên làm việc", "Phần chung")
    steps = ["Mở http://localhost:8077/ → nhập tên đăng nhập & mật khẩu → Đăng nhập",
             "Phiên ghi nhớ 12 giờ — mở lại tab không cần đăng nhập lại",
             "Sai nhiều lần → bị giới hạn tần suất (10/phút/IP), chờ ~1 phút",
             "Lần đầu với admin (mật khẩu mặc định) → buộc đổi mật khẩu ngay"]
    for i, st in enumerate(steps):
        t = 1.85 + i * 0.92
        circle(s, 0.6, t, 0.58, [NAVY, TEAL, SEA, AMBER][i], str(i + 1), WHITE, 18)
        rect(s, 1.4, t, 11.3, 0.72, CARD if i % 2 == 0 else CARD2, rounded=True, line=LINEC, lw=0.6)
        tb(s, 1.65, t, 10.9, 0.72, [one(st, 13.5, INK, line=1.0)], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, "Đổi mật khẩu & sửa họ tên bất kỳ lúc nào trong tab Hồ sơ")

    # 4. Bảng tài khoản
    s = blank(prs); heading(s, "10 tài khoản demo", "Mật khẩu demo: 123456 · admin: admin123")
    rows = [["Tài khoản", "Chức danh", "Vai trò"],
            ["admin", "Quản trị hệ thống", "admin"],
            ["giamdoc", "Giám đốc nhà máy", "supervisor (xem)"],
            ["quandoc", "Quản đốc phân xưởng", "supervisor"],
            ["truongca", "Trưởng ca sản xuất", "supervisor"],
            ["vanhanh", "Nhân viên vận hành", "operator"],
            ["kcs", "Nhân viên KCS / QA", "qa"],
            ["kysu", "Kỹ sư công nghệ", "engineer"],
            ["thukho", "Thủ kho NVL", "operator"],
            ["baotri", "Nhân viên bảo trì", "operator"],
            ["nangluong", "NV quản lý năng lượng", "operator"]]
    simple_table(s, 0.6, 1.7, 9.0, rows, [0.28, 0.46, 0.26], rh=0.345, bsize=11, hsize=11)
    rect(s, 9.9, 1.7, 2.8, 3.95, CARD, rounded=True, line=LINEC, lw=0.75)
    tb(s, 10.1, 1.9, 2.45, 3.6,
       [one("Lưu ý", 14, AMBER, b=True, sa=6),
        one("Đây là tài khoản demo.", 12, INK, sa=6, line=1.1),
        one("Môi trường thật chỉ tạo admin; admin tự tạo tài khoản theo chức danh thực tế.",
            12, GREY, line=1.12)])

    # 5. 4 lớp phân quyền
    s = blank(prs); heading(s, "Vì sao tôi không bấm được nút nào đó?", "4 lớp kiểm soát")
    layers = [("1", "Menu", "Chức danh quyết định tab nào hiện. Tab Hồ sơ luôn có.", NAVY),
              ("2", "Quyền thao tác", "Nút ghi/tạo/duyệt chỉ chạy nếu có quyền — thiếu thì báo 403.", TEAL),
              ("3", "Phạm vi (scope)", "Chỉ thao tác trên line / khu vực / loại test được phân.", SEA),
              ("4", "Phân tách (SoD)", "Không tự duyệt việc mình làm: soạn≠duyệt, ghi QC≠release, ký≠khóa.", AMBER)]
    for i, (bd, hd, ds, col) in enumerate(layers):
        t = 1.8 + i * 1.22
        circle(s, 0.6, t + 0.1, 0.72, col, bd, WHITE, 24)
        rect(s, 1.55, t, 11.15, 1.0, CARD, rounded=True, line=LINEC, lw=0.6)
        tb(s, 1.85, t + 0.12, 3.0, 0.76, [one(hd, 16, NAVY, b=True)], anchor=MSO_ANCHOR.MIDDLE)
        tb(s, 4.7, t + 0.12, 7.8, 0.76, [one(ds, 13, INK, line=1.05)], anchor=MSO_ANCHOR.MIDDLE)

    # 6. Giao diện chung
    s = blank(prs); heading(s, "Làm quen giao diện", "Phần chung")
    ui = [("Thanh menu", "Các tab chức năng — chỉ hiện tab được phân quyền"),
          ("Subnav", "Thanh phụ chia mục con (vd Kho: Tồn / Thẻ kho / Hạn dùng)"),
          ("Ô tìm kiếm", "Gõ để lọc nhanh dòng trên bảng lớn"),
          ("Modal", "Cửa sổ bật lên: tạo/sửa/xem chi tiết (BOM, EBR, tem)"),
          ("Toast", "Thông báo góc màn hình: xanh = thành công, đỏ = lỗi"),
          ("Màu trạng thái", "Đỏ=thiếu thông tin · xanh lá=chưa NVL · xanh dương=đầy đủ")]
    for i, (h, d) in enumerate(ui):
        col = i % 2; row = i // 2
        l = 0.6 + col * 6.15; t = 1.8 + row * 1.55
        rect(s, l, t, 5.95, 1.35, CARD, rounded=True, line=LINEC, lw=0.7)
        tb(s, l + 0.25, t + 0.16, 5.5, 0.45, [one(h, 14.5, TEAL, b=True)])
        tb(s, l + 0.25, t + 0.62, 5.5, 0.65, [one(d, 11.5, INK, line=1.08)])

    # 7-10. Theo chức danh
    s = blank(prs); heading(s, "Quản đốc & Trưởng ca", "Điều hành sản xuất")
    card(s, 0.6, 1.7, 5.95, 4.8, "Q", "Quản đốc phân xưởng",
         ["Tạo lệnh SX & Work Order", "Dispatch + lập lịch tối ưu",
          "Theo dõi tiến độ mẻ", "Mở/đóng deviation",
          "Phê duyệt & KHÓA hồ sơ EBR", "(toàn nhà máy)"], NAVY, bsize=13)
    card(s, 6.75, 1.7, 5.95, 4.8, "T", "Trưởng ca sản xuất",
         ["Tạo mẻ (kiểm tra tồn trước)", "Chuyển trạng thái ready→running",
          "Consume / produce / actual", "Chạy phase ISA-88",
          "Ký điện tử EBR", "Phạm vi: Line Nấu A"], TEAL, bsize=13)

    s = blank(prs); heading(s, "Vận hành & KCS", "Thao tác sàn & chất lượng")
    card(s, 0.6, 1.7, 5.95, 4.8, "V", "Nhân viên vận hành",
         ["Chọn mẻ đang chạy", "Consume NVL theo định mức",
          "Ghi actual tham số", "Produce lô output",
          "Chạy phase ISA-88", "Khuyên dùng giao diện Kiosk"], SEA, bsize=13)
    card(s, 6.75, 1.7, 5.95, 4.8, "K", "Nhân viên KCS / QA",
         ["Ghi kết quả QC (FAIL → hold)", "RELEASE (chặn nếu còn FAIL)",
          "Mở/xử lý deviation", "SPC · CAPA · COA · LIMS",
          "Duyệt công thức (≠ tự soạn)", "Phạm vi test: Độ đường, pH"], NAVY, bsize=13)

    s = blank(prs); heading(s, "Kỹ sư & Thủ kho", "Thiết kế & kho")
    card(s, 0.6, 1.7, 5.95, 4.8, "E", "Kỹ sư công nghệ",
         ["Tạo sản phẩm/vật tư", "Soạn công thức + BOM + ISA-88",
          "→ review → duyệt → hiệu lực", "(không tự duyệt bản mình soạn)",
          "Yield, change-control, NVL", "   thay thế"], TEAL, bsize=13)
    card(s, 6.75, 1.7, 5.95, 4.8, "W", "Thủ kho NVL",
         ["Nhập/xuất/hoàn/sang ngang", "Tồn · thẻ kho · hạn dùng",
          "Đóng pallet thành phẩm (WMS)", "Putaway / ship + in tem",
          "Ghi biến động bao bì", "Phạm vi: khu vực kho"], NAVY, bsize=13)

    s = blank(prs); heading(s, "Bảo trì · Năng lượng · Giám đốc · Admin", "Các chức danh khác")
    card(s, 0.6, 1.7, 3.86, 4.8, "🔧", "Bảo trì",
         ["Báo cáo + xử lý sự cố", "Kế hoạch bảo trì",
          "Kiểm định/hiệu chuẩn", "Phụ tùng tồn min",
          "KV: lọc, chiết"], TEAL, bsize=12.5, hsize=15)
    card(s, 4.73, 1.7, 3.86, 4.8, "⚡", "Năng lượng",
         ["Cập nhật số liệu ngày", "Biểu đồ ngày/tháng",
          "Quản danh mục nhóm/khu", "KV: nấu, lên men, chiết"], SEA, bsize=12.5, hsize=15)
    card(s, 8.86, 1.7, 3.86, 4.8, "★", "Giám đốc & Admin",
         ["Giám đốc: chỉ xem —", "   KPI, OEE, BC, truy xuất, AI",
          "Admin: quản lý tài khoản,", "   quyền, scope, API key,",
          "   webhook, audit"], NAVY, bsize=12.5, hsize=15)

    # 11. Quy trình end-to-end
    s = blank(prs); heading(s, "Quy trình end-to-end một mẻ", "Ai làm gì")
    rows = [["Bước", "Người", "Thao tác"],
            ["1-2", "Quản đốc / Trưởng ca", "Tạo lệnh SX → Work Order → phát hành"],
            ["3", "Thủ kho", "Đảm bảo NVL đã nhập kho, còn hạn"],
            ["4-5", "Trưởng ca / Vận hành", "Kiểm tra tồn → tạo mẻ → cấp liệu (FEFO)"],
            ["6", "Vận hành", "Chạy phase · ghi actual · consume"],
            ["7", "KCS", "Ghi QC; xử lý deviation nếu FAIL"],
            ["8-9", "Vận hành / KCS", "Produce lô → RELEASE (hết FAIL)"],
            ["10-12", "Vận hành / Quản đốc", "completed → ký EBR → khóa → closed"],
            ["13", "Thủ kho", "Đóng pallet, putaway, in tem"]]
    simple_table(s, 0.6, 1.7, 12.1, rows, [0.12, 0.30, 0.58], rh=0.43, bsize=11)

    # 12. Kiosk
    s = blank(prs); heading(s, "Kiosk xưởng", "/kiosk.html — màn hình cảm ứng")
    card(s, 0.6, 1.85, 2.95, 4.4, "📷", "Quét mã",
         ["Bắn mã / gõ + Enter", "Tự phân giải lô/mẻ/", "   WO/đơn", "Cấp liệu nhanh nút", "   SL lớn"], TEAL, hsize=15)
    card(s, 3.78, 1.85, 2.95, 4.4, "🏷️", "In tem",
         ["Nhập mã → Tạo tem", "Code39", "In trực tiếp"], SEA, hsize=15)
    card(s, 6.96, 1.85, 2.95, 4.4, "⚗️", "Mẻ chạy",
         ["Xem nhanh các mẻ", "đang chạy"], NAVY, hsize=15)
    card(s, 10.14, 1.85, 2.58, 4.4, "💻", "Bản đầy đủ",
         ["Quay lại UI", "chính"], AMBER, hsize=15)
    footer(s, "Tối ưu cho tablet/máy quét tại sàn — đăng nhập bằng tài khoản của bạn")

    # 13. Trợ lý AI
    s = blank(prs); heading(s, "Trợ lý AI", "Hỏi-đáp + cảnh báo vận hành")
    card(s, 0.6, 1.85, 5.95, 4.2, "✦", "Khung chat (trái)",
         ["Tạo / chọn hội thoại (lưu server)", "Gõ câu hỏi → Gửi (hiện dần)",
          "Hiện tool AI đã dùng", "Còn nguyên khi tải lại / đổi máy"], TEAL, bsize=13, hsize=15)
    card(s, 6.75, 1.85, 5.95, 4.2, "▲", "Cảnh báo (phải)",
         ["3 mức: Cao / TB / Thấp", "Miền · phát hiện · đề xuất",
          "Nút Báo cáo nền (job async)", "AI chỉ tư vấn — không tự đổi dữ liệu"], NAVY, bsize=13, hsize=15)
    footer(s, "Không có API key → engine luật offline · có hạn mức chat/ngày để kiểm soát chi phí")

    # 14. FAQ
    s = blank(prs); heading(s, "Xử lý sự cố thường gặp", "FAQ")
    rows = [["Tình huống", "Cách xử lý"],
            ["Không thấy một tab", "Chức danh không được cấp menu — nhờ admin"],
            ["Bấm nút → 'không đủ quyền' (403)", "Thiếu quyền thao tác — nhờ admin gán"],
            ["Không thấy mẻ line khác", "Phạm vi (scope) giới hạn — đúng thiết kế"],
            ["Không duyệt/khóa được việc mình làm", "Phân tách nhiệm vụ (SoD) — nhờ người khác"],
            ["Tạo mẻ báo 'thiếu tồn' (409)", "Nhập kho thêm hoặc xác nhận cho phép thiếu"],
            ["Release bị chặn", "Đóng deviation (hết FAIL) rồi release"],
            ["Mẻ không sửa được nữa", "EBR đã khóa → mẻ bất biến, chỉ amendment"],
            ["AI báo hết lượt", "Vượt hạn mức/ngày — chờ hoặc tăng quota"]]
    simple_table(s, 0.6, 1.7, 12.1, rows, [0.42, 0.58], rh=0.43, bsize=11)

    # 15. Closing
    s = blank(prs, NAVY)
    circle(s, 5.3, 1.7, 1.5, TEAL, "🍺", WHITE, 34)
    tb(s, 1.0, 3.4, 11.3, 1.0,
       [one("Sẵn sàng vận hành MES", 34, WHITE, b=True, align=PP_ALIGN.CENTER, font=HEAD)])
    tb(s, 1.0, 4.5, 11.3, 0.7,
       [one("Đọc chi tiết trong Sách Hướng dẫn Sử dụng · Tài liệu API tại /docs",
            16, ICE, align=PP_ALIGN.CENTER)])
    tb(s, 1.0, 6.7, 11.3, 0.4,
       [one("MES Bia Hạ Long · Nhà Máy Đông Mai · 0.1.0-mvp", 12, ICE, align=PP_ALIGN.CENTER)])

    prs.save("docs/MES-HuongDanSuDung.pptx")
    print("✓ docs/MES-HuongDanSuDung.pptx")


if __name__ == "__main__":
    deck_architecture()
    deck_features()
    deck_userguide()
